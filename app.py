import streamlit as st
import os
import json
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Set page config
st.set_page_config(
    page_title="Pinokio Presenter",
    page_icon="🎬",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Create data directory if it doesn't exist
DATA_DIR = Path("presentations")
DATA_DIR.mkdir(exist_ok=True)

# Custom CSS
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        color: #FF6B6B;
        margin-bottom: 1rem;
    }
    .slide-preview {
        border: 2px solid #FF6B6B;
        border-radius: 8px;
        padding: 1rem;
        margin: 0.5rem 0;
        background-color: #f8f9fa;
    }
    .success-message {
        color: #28a745;
        font-weight: bold;
    }
</style>
""", unsafe_allow_html=True)

# Title
st.markdown('<div class="main-header">🎬 Pinokio Presenter</div>', unsafe_allow_html=True)
st.markdown("Create and display beautiful presentations for Pinokio AI")

# Sidebar
st.sidebar.title("🎯 Navigation")
mode = st.sidebar.radio("Select Mode", ["Create Presentation", "View Presentations", "Edit Presentation"])

def save_presentation_data(name, slides):
    """Save presentation data to JSON"""
    filepath = DATA_DIR / f"{name}.json"
    with open(filepath, 'w') as f:
        json.dump({
            "name": name,
            "created": datetime.now().isoformat(),
            "slides": slides
        }, f, indent=2)
    return filepath

def load_presentation_data(name):
    """Load presentation data from JSON"""
    filepath = DATA_DIR / f"{name}.json"
    if filepath.exists():
        with open(filepath, 'r') as f:
            return json.load(f)
    return None

def export_to_pptx(presentation_data):
    """Export presentation to PowerPoint format"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for slide_data in presentation_data["slides"]:
        # Add blank slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Add title
        if slide_data.get("title"):
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = slide_data["title"]
            title_frame.paragraphs[0].font.size = Pt(54)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 107, 107)
        
        # Add content
        if slide_data.get("content"):
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(9)
            height = Inches(5)
            content_box = slide.shapes.add_textbox(left, top, width, height)
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.text = slide_data["content"]
            content_frame.paragraphs[0].font.size = Pt(24)
    
    return prs

# CREATE PRESENTATION MODE
if mode == "Create Presentation":
    st.header("✨ Create New Presentation")
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        presentation_name = st.text_input(
            "Presentation Title",
            placeholder="e.g., AI for Beginners",
            key="pres_title"
        )
    
    with col2:
        num_slides = st.number_input("Number of Slides", min_value=1, max_value=20, value=3)
    
    if presentation_name:
        st.markdown("---")
        st.subheader("📝 Slide Content")
        
        slides = []
        
        for i in range(num_slides):
            with st.container():
                st.markdown(f"**Slide {i+1}**")
                
                col1, col2 = st.columns(2)
                
                with col1:
                    slide_title = st.text_input(
                        f"Slide {i+1} Title",
                        placeholder="Enter slide title",
                        key=f"title_{i}"
                    )
                
                with col2:
                    slide_type = st.selectbox(
                        f"Slide {i+1} Type",
                        ["Title Slide", "Content Slide", "Quote Slide"],
                        key=f"type_{i}"
                    )
                
                slide_content = st.text_area(
                    f"Slide {i+1} Content",
                    placeholder="Enter slide content (bullet points or paragraphs)",
                    height=100,
                    key=f"content_{i}"
                )
                
                slides.append({
                    "title": slide_title,
                    "content": slide_content,
                    "type": slide_type
                })
                
                st.markdown("---")
        
        # Create buttons
        col1, col2, col3 = st.columns(3)
        
        with col1:
            if st.button("💾 Save Presentation", use_container_width=True):
                if presentation_name.strip():
                    save_presentation_data(presentation_name, slides)
                    st.success(f"✅ Presentation '{presentation_name}' saved successfully!")
                else:
                    st.error("Please enter a presentation name")
        
        with col2:
            if st.button("📊 Preview Slides", use_container_width=True):
                st.info("Preview below ↓")
        
        with col3:
            if st.button("📥 Export to PowerPoint", use_container_width=True):
                if presentation_name.strip():
                    prs = export_to_pptx({"slides": slides})
                    output_path = DATA_DIR / f"{presentation_name}.pptx"
                    prs.save(str(output_path))
                    st.success(f"✅ Exported to {output_path}")
                else:
                    st.error("Please enter a presentation name")

# VIEW PRESENTATIONS MODE
elif mode == "View Presentations":
    st.header("📚 Your Presentations")
    
    # List saved presentations
    json_files = list(DATA_DIR.glob("*.json"))
    
    if json_files:
        selected_pres = st.selectbox(
            "Select a presentation to view",
            [f.stem for f in json_files]
        )
        
        if selected_pres:
            pres_data = load_presentation_data(selected_pres)
            
            if pres_data:
                st.markdown(f"## {pres_data['name']}")
                st.markdown(f"*Created: {pres_data['created'][:10]}*")
                
                # Slide navigation
                st.markdown("---")
                slide_num = st.slider(
                    "Select Slide",
                    1,
                    len(pres_data["slides"]),
                    1
                )
                
                # Display current slide
                current_slide = pres_data["slides"][slide_num - 1]
                
                st.markdown(f"### {current_slide['title']}")
                st.markdown(f"> **Type:** {current_slide['type']}")
                st.markdown("---")
                st.markdown(current_slide['content'])
                
                # Export option
                col1, col2 = st.columns(2)
                with col1:
                    if st.button("📥 Export This Presentation to PowerPoint"):
                        prs = export_to_pptx(pres_data)
                        output_path = DATA_DIR / f"{selected_pres}_export.pptx"
                        prs.save(str(output_path))
                        st.success(f"✅ Exported to {output_path}")
                
                with col2:
                    if st.button("🗑️ Delete Presentation"):
                        json_file = DATA_DIR / f"{selected_pres}.json"
                        json_file.unlink()
                        st.success(f"Presentation '{selected_pres}' deleted")
                        st.rerun()
    else:
        st.info("📭 No presentations found. Create one in the 'Create Presentation' tab!")

# EDIT PRESENTATION MODE
elif mode == "Edit Presentation":
    st.header("✏️ Edit Presentation")
    
    json_files = list(DATA_DIR.glob("*.json"))
    
    if json_files:
        selected_pres = st.selectbox(
            "Select a presentation to edit",
            [f.stem for f in json_files],
            key="edit_select"
        )
        
        if selected_pres:
            pres_data = load_presentation_data(selected_pres)
            
            if pres_data:
                st.markdown(f"### Editing: {pres_data['name']}")
                
                slides = pres_data["slides"]
                
                for i in range(len(slides)):
                    with st.container():
                        st.markdown(f"**Slide {i+1}**")
                        
                        col1, col2 = st.columns(2)
                        
                        with col1:
                            slides[i]["title"] = st.text_input(
                                f"Title for Slide {i+1}",
                                value=slides[i].get("title", ""),
                                key=f"edit_title_{i}"
                            )
                        
                        with col2:
                            slides[i]["type"] = st.selectbox(
                                f"Type for Slide {i+1}",
                                ["Title Slide", "Content Slide", "Quote Slide"],
                                index=["Title Slide", "Content Slide", "Quote Slide"].index(slides[i].get("type", "Content Slide")),
                                key=f"edit_type_{i}"
                            )
                        
                        slides[i]["content"] = st.text_area(
                            f"Content for Slide {i+1}",
                            value=slides[i].get("content", ""),
                            height=100,
                            key=f"edit_content_{i}"
                        )
                        
                        st.markdown("---")
                
                if st.button("💾 Save Changes", use_container_width=True):
                    save_presentation_data(selected_pres, slides)
                    st.success(f"✅ Changes saved to '{selected_pres}'!")
    else:
        st.info("📭 No presentations found to edit.")

# Footer
st.markdown("---")
st.markdown("""
<div style='text-align: center; color: #666; font-size: 0.9rem;'>
    Made for Pinokio AI | Create amazing presentations locally
</div>
""", unsafe_allow_html=True)